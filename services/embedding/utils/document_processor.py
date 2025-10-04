import fitz  # PyMuPDF
import docx
from pptx import Presentation
import io
from typing import List, Dict, Any
import re

class DocumentProcessor:
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.pptx']
    
    def extract_text_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF file content using PyMuPDF."""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_document = fitz.open(stream=pdf_file, filetype="pdf")
            text = ""
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                page_text = page.get_text()
                if page_text.strip():  # Only add non-empty pages
                    text += f"--- Page {page_num + 1} ---\n{page_text}\n\n"
            
            pdf_document.close()
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting PDF text: {str(e)}")
    
    def extract_text_with_page_info(self, file_content: bytes, filename: str) -> tuple[str, list]:
        """Extract text from PDF with page information for better chunking."""
        try:
            file_extension = filename.lower().split('.')[-1]
            
            if file_extension == 'pdf':
                pdf_file = io.BytesIO(file_content)
                pdf_document = fitz.open(stream=pdf_file, filetype="pdf")
                text = ""
                page_info = []
                
                for page_num in range(pdf_document.page_count):
                    page = pdf_document[page_num]
                    page_text = page.get_text()
                    if page_text.strip():
                        text += page_text + "\n"
                        page_info.append({
                            'page_number': page_num + 1,
                            'start_char': len(text) - len(page_text),
                            'end_char': len(text),
                            'text_length': len(page_text)
                        })
                
                pdf_document.close()
                return text.strip(), page_info
            else:
                # For non-PDF files, return text without page info
                text = self.extract_text(file_content, filename)
                return text, []
                
        except Exception as e:
            raise Exception(f"Error extracting text with page info: {str(e)}")
    
    def extract_text_from_docx(self, file_content: bytes) -> str:
        """Extract text from DOCX file content."""
        try:
            doc_file = io.BytesIO(file_content)
            doc = docx.Document(doc_file)
            text = ""
            
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting DOCX text: {str(e)}")
    
    def extract_text_from_pptx(self, file_content: bytes) -> str:
        """Extract text from PPTX file content."""
        try:
            ppt_file = io.BytesIO(file_content)
            prs = Presentation(ppt_file)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                text += f"--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting PPTX text: {str(e)}")
    
    def extract_text(self, file_content: bytes, filename: str) -> str:
        """Extract text from supported file formats."""
        file_extension = filename.lower().split('.')[-1]
        
        if file_extension == 'pdf':
            return self.extract_text_from_pdf(file_content)
        elif file_extension == 'docx':
            return self.extract_text_from_docx(file_content)
        elif file_extension == 'pptx':
            return self.extract_text_from_pptx(file_content)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    def clean_text(self, text: str) -> str:
        """Clean and normalize extracted text."""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove special characters but keep basic punctuation
        text = re.sub(r'[^\w\s.,!?;:()-]', '', text)
        return text.strip()
    
    def is_supported_format(self, filename: str) -> bool:
        """Check if file format is supported."""
        file_extension = '.' + filename.lower().split('.')[-1]
        return file_extension in self.supported_formats
